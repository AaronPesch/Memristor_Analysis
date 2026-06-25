import io
import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

import PySide6.QtWidgets as qt
from PySide6.QtCore import QThread, Qt, QUrl
from PySide6.QtGui import QDesktopServices
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas as rl_canvas

from ..converter import BatchConverter, path_to_glob
from ..core import MenuAction, Mode
from .import_worker import ImportWorker
from .menu_bar import MenuBar
from .navigation_bar import NavigationBar


class MainWindow(qt.QMainWindow):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Memristor Analysis Tool")
        self.resize(1200, 800)
        self.setup_ui()
        self.setup_connections()

    def setup_ui(self) -> None:
        self.menu_bar = MenuBar(self)
        self.setMenuBar(self.menu_bar)
        central_widget = qt.QWidget()
        self.setCentralWidget(central_widget)
        self.main_layout = qt.QVBoxLayout(central_widget)
        self.nav_bar = NavigationBar()
        self.main_layout.addWidget(self.nav_bar)

    def setup_connections(self):
        menu_actions = self.menu_bar.menu_actions

        menu_actions[MenuAction.EXIT].triggered.connect(self.cleanup_and_exit)

        # Connect Imports (Shortcuts and Menu)
        menu_actions[MenuAction.IMPORT_DEVICE].triggered.connect(
            lambda: self.handle_import(mode=Mode.DEVICE)
        )
        menu_actions[MenuAction.IMPORT_STACK].triggered.connect(
            lambda: self.handle_import(mode=Mode.STACK)
        )

        menu_actions[MenuAction.EXPORT_ALL].triggered.connect(
            lambda checked=False: self.export_all("png")
        )
        menu_actions[MenuAction.EXPORT_ALL_PDF_COMBINED].triggered.connect(
            lambda checked=False: self.export_all_combined_pdf()
        )
        menu_actions[MenuAction.EXPORT_ALL_PPTX].triggered.connect(
            lambda checked=False: self.export_all_combined_pptx()
        )
        menu_actions[MenuAction.EXPORT_CURRENT].triggered.connect(
            lambda checked=False: self.export_current("png")
        )

        # Export menu for current
        menu_actions[MenuAction.EXPORT_CURRENT_PNG].triggered.connect(
            lambda checked=False: self.export_current("png")
        )

        menu_actions[MenuAction.EXPORT_CURRENT_JPEG].triggered.connect(
            lambda checked=False: self.export_current("jpeg")
        )

        menu_actions[MenuAction.EXPORT_CURRENT_EPS].triggered.connect(
            lambda checked=False: self.export_current("eps")
        )
        menu_actions[MenuAction.EXPORT_CURRENT_SVG].triggered.connect(
            lambda checked=False: self.export_current("SVG")
        )
        menu_actions[MenuAction.EXPORT_CURRENT_PDF].triggered.connect(
            lambda checked=False: self.export_current("pdf")
        )
        menu_actions[MenuAction.EXPORT_CURRENT_CSV].triggered.connect(
            lambda checked=False: self.export_current("csv")
        )
        menu_actions[MenuAction.EXPORT_CURRENT_TXT].triggered.connect(
            lambda checked=False: self.export_current("txt")
        )

        # Export menu for all

        menu_actions[MenuAction.EXPORT_ALL_PNG].triggered.connect(
            lambda checked=False: self.export_all("png")
        )

        menu_actions[MenuAction.EXPORT_ALL_JPEG].triggered.connect(
            lambda checked=False: self.export_all("jpeg")
        )

        menu_actions[MenuAction.EXPORT_ALL_EPS].triggered.connect(
            lambda checked=False: self.export_all("eps")
        )
        menu_actions[MenuAction.EXPORT_ALL_SVG].triggered.connect(
            lambda checked=False: self.export_all("svg")
        )
        menu_actions[MenuAction.EXPORT_ALL_PDF].triggered.connect(
            lambda checked=False: self.export_all("pdf")
        )
        menu_actions[MenuAction.EXPORT_ALL_CSV].triggered.connect(
            lambda checked=False: self.export_all("csv")
        )
        menu_actions[MenuAction.EXPORT_ALL_TXT].triggered.connect(
            lambda checked=False: self.export_all("txt")
        )

        # Link to GH-Wiki
        menu_actions[MenuAction.VIEW_HELP].triggered.connect(self.open_wiki)

    def export_current(self, fmt: str):

        viewer = self.nav_bar.get_current_viewer()

        if viewer is None:
            qt.QMessageBox.warning(self, "Export", "No plot selected")
            return

        file_path, _ = qt.QFileDialog.getSaveFileName(
            self,
            "Export Plot",
            f"plot.{fmt}",
            f"{fmt.upper()} Files (*.{fmt})",
        )

        if not file_path:
            return

        if fmt in {"csv", "txt"}:
            ok = viewer.export_data(file_path, fmt)
        else:
            ok = viewer.export_image(file_path, fmt)

        if not ok:
            qt.QMessageBox.warning(
                self,
                "Export",
                "Export failed: missing .json next to the loaded .html.",
            )

    def export_all_combined_pptx(self):
        viewers = self.nav_bar.get_all_viewers()
        if not viewers:
            qt.QMessageBox.warning(self, "Export", "No plots loaded")
            return

        file_path, _ = qt.QFileDialog.getSaveFileName(
            self, "Export PowerPoint", "all_plots.pptx", "PowerPoint Files (*.pptx)"
        )
        if not file_path:
            return

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        blank_layout = prs.slide_layouts[6]
        exported = 0

        for viewer in viewers:
            fig = viewer._resolve_figure()
            if fig is None:
                continue
            try:
                fig = viewer._prepare_for_export(fig)
                img_bytes = fig.to_image(format="png", scale=2)
                slide = prs.slides.add_slide(blank_layout)
                slide.shapes.add_picture(
                    io.BytesIO(img_bytes),
                    0,
                    0,
                    width=prs.slide_width,
                    height=prs.slide_height,
                )
                exported += 1
            except Exception as e:
                print(f"Skipping plot: {e}")

        if exported == 0:
            qt.QMessageBox.warning(self, "Export", "No plots could be exported.")
            return

        prs.save(file_path)

    def export_all_combined_pdf(self):
        viewers = self.nav_bar.get_all_viewers()
        if not viewers:
            qt.QMessageBox.warning(self, "Export", "No plots loaded")
            return

        file_path, _ = qt.QFileDialog.getSaveFileName(
            self, "Export Combined PDF", "all_plots.pdf", "PDF Files (*.pdf)"
        )
        if not file_path:
            return

        c = rl_canvas.Canvas(file_path)
        exported = 0

        for viewer in viewers:
            fig = viewer._resolve_figure()
            if fig is None:
                continue
            try:
                fig = viewer._prepare_for_export(fig)
                img_bytes = fig.to_image(format="png", scale=2)
                img_reader = ImageReader(io.BytesIO(img_bytes))
                w, h = img_reader.getSize()
                c.setPageSize((w, h))
                c.drawImage(img_reader, 0, 0, w, h)
                c.showPage()
                exported += 1
            except Exception as e:
                print(f"Skipping plot: {e}")

        if exported == 0:
            qt.QMessageBox.warning(self, "Export", "No plots could be exported.")
            return

        c.save()

    def export_all(self, fmt: str):

        folder = qt.QFileDialog.getExistingDirectory(self, "Select Export Folder")

        if not folder:
            return

        viewers = self.nav_bar.get_all_viewers()

        if not viewers:
            qt.QMessageBox.warning(self, "Export", "No plots loaded")
            return

        for i, viewer in enumerate(viewers):
            if getattr(viewer, "html_path", None):
                html = Path(viewer.html_path)
                filename = f"{html.parent.name}_{html.stem}"
            else:
                filename = f"plot_{i}"

            path = Path(folder) / f"{filename}.{fmt}"

            if fmt in {"csv", "txt"}:
                viewer.export_data(str(path), fmt)
            else:
                viewer.export_image(str(path), fmt)

    # Helper to get the figure
    def _get_figure(self, viewer):

        if hasattr(viewer, "get_figure"):
            return viewer.get_figure()

        if hasattr(viewer, "figure"):
            return viewer.figure

        return None

    # Helper to save the figure
    def _write_figure(self, fig, path):

        path = str(path)

        # Plotly
        if hasattr(fig, "write_image"):
            fig.write_image(path)
            return

        # Matplotlib
        if hasattr(fig, "savefig"):
            fig.savefig(path)
            return

        qt.QMessageBox.warning(self, "Export", "Unsupported figure type")

    def handle_import(self, mode: Mode):
        folder = qt.QFileDialog.getExistingDirectory(
            self, f"Select {mode.value} Folder"
        )
        if not folder:
            return

        path = path_to_glob(folder, mode)

        # 1. Create Progress Dialog
        self.pd = qt.QProgressDialog("Initializing...", None, 0, 100, self)
        self.pd.setWindowTitle("Processing Data")
        self.pd.setWindowModality(Qt.WindowModal)
        self.pd.setMinimumDuration(0)
        self.pd.show()

        # 2. Setup Thread and Worker
        self.import_thread = QThread()
        self.worker = ImportWorker(path, mode, BatchConverter)
        self.worker.moveToThread(self.import_thread)

        # 3. Connect Signals
        self.import_thread.started.connect(self.worker.run)
        self.worker.progress.connect(self.pd.setValue)
        self.worker.status_message.connect(self.pd.setLabelText)

        # Clean up thread when finished
        self.worker.finished.connect(self.import_thread.quit)
        self.worker.finished.connect(self.worker.deleteLater)
        self.import_thread.finished.connect(self.import_thread.deleteLater)

        # UI Refresh on success
        self.worker.finished.connect(self.on_import_success)
        self.worker.error.connect(self.on_import_error)

        # 4. Start
        self.import_thread.start()

    def on_import_success(self):
        self.pd.close()
        # Force the NavigationBar to reload the tabs (which now have new HTML files)
        self.nav_bar.update_tabs_by_level()
        qt.QMessageBox.information(
            self, "Success", "Data imported and plots generated successfully."
        )

    def on_import_error(self, err_msg):
        self.pd.close()
        qt.QMessageBox.critical(
            self, "Error", f"An error occurred during processing:\n{err_msg}"
        )

    def apply_to_active(self, callback):
        viewer = self.nav_bar.get_current_viewer()
        if viewer:
            callback(viewer)

    def cleanup_and_exit(self):
        """Deletes contents of src/app/temp/ and exits the application."""
        temp_dir = Path(__file__).parent.parent / "temp"
        print(f"Cleaning up temporary files in {temp_dir}...")
        if temp_dir.exists() and temp_dir.is_dir():
            shutil.rmtree(temp_dir)
        self.close()

    def open_wiki(self):
        """Opens the GitHub Wiki in the default browser."""
        url = QUrl(
            "https://github.com/arrayan/memristor-analysis/wiki/Memristor-Analysis-Tool-Wiki"
        )
        QDesktopServices.openUrl(url)
